import sys
import subprocess

# Install needed libraries quietly
try:
    import pptx
    import docx
except ImportError:
    subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "python-docx", "-q"])
    import pptx
    import docx

import os

with open(r"d:\RagaVoiceStudio\extracted_docs.txt", "w", encoding="utf-8") as f:
    f.write("=== PPTX CONTENT ===\n")
    try:
        prs = pptx.Presentation(r"d:\RagaVoiceStudio\STATIC\REVIEW-1-PPT.pptx")
        for i, slide in enumerate(prs.slides):
            f.write(f"\n--- Slide {i+1} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text.replace('\n', ' ') + "\n")
    except Exception as e:
        f.write(f"Error reading PPTX: {e}\n")

    f.write("\n=== DOCX STRUCTURAL OUTLINE ===\n")
    try:
        doc = docx.Document(r"d:\RagaVoiceStudio\STATIC\final_thesis_dot.docx")
        for p in doc.paragraphs:
            if p.text.strip():
                f.write(p.text + "\n")
    except Exception as e:
        f.write(f"Error reading DOCX: {e}\n")

